"""
Document Manager Module with Supabase Integration
Handles document upload, processing, and management with Supabase storage
"""

import os
import logging
import asyncio
import tempfile
from typing import List, Dict, Any, Optional, Union
from pathlib import Path
from datetime import datetime
import mimetypes
import aiofiles
import PyPDF2
from docx import Document
from pptx import Presentation
import openpyxl
import markdown
import json
import uuid

from .retrieval_client import RetrievalAPIClient
from .supabase_client import get_supabase_manager, reset_supabase_manager
from .utils import extract_file_metadata, sanitize_filename, parse_file_type, format_file_size
from .usage_tracker import get_usage_tracker


logger = logging.getLogger(__name__)


class DocumentManagerSupabase:
    """Manages document uploads and processing with Supabase integration"""
    
    SUPPORTED_EXTENSIONS = {
        '.pdf', '.doc', '.docx', '.txt', '.md', 
        '.ppt', '.pptx', '.xls', '.xlsx', '.csv',
        '.json', '.xml', '.html', '.htm'
    }
    
    MAX_FILE_SIZE = 30 * 1024 * 1024  # 30 MB (Cloud Run request size limit)
    STORAGE_BUCKET = "documents"
    
    def __init__(self, retrieval_client: RetrievalAPIClient):
        self.retrieval_client = retrieval_client
        # TEMPORARY: Force reload to pick up corrected env vars
        reset_supabase_manager()
        self.supabase = get_supabase_manager()
        self.usage_tracker = get_usage_tracker()
        
    async def upload_document(self, file_content: bytes = None, filename: str = None, 
                            content_type: str = None, file_path: str = None, 
                            metadata: Optional[Dict[str, Any]] = None,
                            user_id: Optional[str] = None,
                            uploaded_by_email: Optional[str] = None) -> Dict[str, Any]:
        """Upload a document to Supabase and OpenAI Retrieval API"""
        logger.info(f"[UPLOAD START] Starting document upload")
        logger.info(f"  - Filename: {filename}")
        logger.info(f"  - Content-Type: {content_type}")
        logger.info(f"  - User: {uploaded_by_email} (ID: {user_id})")
        logger.info(f"  - File size: {len(file_content) if file_content else 'from file path'}")
        
        try:
            # Handle file content from API
            if file_content and filename:
                # Generate unique filename
                safe_filename = sanitize_filename(filename)
                timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                unique_filename = f"{timestamp}_{safe_filename}"
                
                # Create metadata
                file_metadata = {
                    "filename": filename,
                    "content_type": content_type or mimetypes.guess_type(filename)[0],
                    "size": len(file_content)
                }
                
                # Save temporarily for OpenAI upload using platform-independent temp directory
                import tempfile
                
                # Try multiple temp directory options
                temp_dir = None
                temp_path = None
                
                # Create temp directory if not exists
                temp_dir = os.path.join(os.getcwd(), 'temp')
                os.makedirs(temp_dir, exist_ok=True)
                
                # Use original filename for temp file (important for OpenAI)
                # Preserve the original extension
                file_ext = os.path.splitext(filename)[1]
                safe_filename = sanitize_filename(filename)
                temp_filename = f"{datetime.now().strftime('%Y%m%d_%H%M%S')}_{safe_filename}"
                temp_path = os.path.join(temp_dir, temp_filename)
                
                try:
                    # Write with original filename structure
                    with open(temp_path, 'wb') as f:
                        f.write(file_content)
                    logger.debug(f"[TEMP FILE] Created temporary file: {temp_path}")
                    logger.debug(f"[TEMP FILE] Original filename preserved: {temp_filename}")
                except Exception as e:
                    logger.error(f"[TEMP FILE] Failed to create temp file: {e}")
                    return {
                        "status": "error",
                        "error": f"Failed to create temporary file: {str(e)}"
                    }
                
                file_path = temp_path
                logger.debug(f"[TEMP FILE] Temporary file ready for upload: {file_path}")
            else:
                # Validate file path
                validation = self._validate_file(file_path)
                if not validation["valid"]:
                    return {
                        "status": "error",
                        "error": validation["error"]
                    }
                
                # Extract metadata
                file_metadata = extract_file_metadata(file_path)
                # Read file content for Supabase upload
                async with aiofiles.open(file_path, 'rb') as f:
                    file_content = await f.read()
                unique_filename = f"{datetime.now().strftime('%Y%m%d_%H%M%S')}_{sanitize_filename(file_metadata['filename'])}"
            
            if metadata:
                file_metadata.update(metadata)
            
            # Extract text preview
            text_preview = await self._extract_text_preview(file_path)
            file_metadata["text_preview"] = text_preview
            
            # Upload to Supabase Storage
            storage_path = f"{unique_filename}"  # Don't include 'documents/' prefix
            logger.info(f"Starting Supabase storage upload for: {storage_path}")
            logger.info(f"Original filename: {file_metadata.get('filename', 'unknown')}")
            logger.info(f"Sanitized storage path: {storage_path}")
            
            try:
                self.supabase.upload_file(
                    self.STORAGE_BUCKET,
                    storage_path,
                    file_content,
                    file_metadata.get("content_type", "application/octet-stream")
                )
                logger.info(f"Successfully uploaded to Supabase storage: {storage_path}")
            except Exception as storage_error:
                logger.error(f"Supabase storage upload failed: {storage_error}")
                logger.error(f"Bucket: {self.STORAGE_BUCKET}")
                logger.error(f"Path: {storage_path}")
                logger.error(f"Content size: {len(file_content)} bytes")
                # Re-raise with more context
                raise Exception(f"Storage upload failed: {storage_error}")
            
            # Upload to OpenAI with original filename
            file_id = None
            try:
                # Pass the original filename to preserve it in OpenAI
                original_filename = file_metadata["filename"]
                
                # Check if it's a scanned PDF or image that needs OCR
                if text_preview and len(text_preview.strip()) < 100 and file_metadata.get("content_type", "").startswith(("image/", "application/pdf")):
                    logger.info(f"[OCR] Document appears to be scanned, creating text version...")
                    
                    # Create a text file with OCR content
                    ocr_filename = f"{os.path.splitext(original_filename)[0]}_OCR.txt"
                    ocr_path = os.path.join(tempfile.gettempdir(), ocr_filename)
                    
                    # Get full OCR text
                    try:
                        from .vision_processor import get_vision_processor
                        vision = get_vision_processor()
                        
                        if file_path.lower().endswith('.pdf'):
                            ocr_text = await vision.extract_text_from_pdf(file_path, max_pages=10)
                        else:
                            ocr_text = await vision.extract_text_from_image(file_path)
                        
                        if ocr_text and len(ocr_text.strip()) > 50:
                            # Save OCR text to file
                            with open(ocr_path, 'w', encoding='utf-8') as f:
                                f.write(f"[OCR 추출 텍스트 - {original_filename}]\n\n")
                                f.write(ocr_text)
                            
                            # Upload OCR text file
                            logger.info(f"[OCR] Uploading OCR text file: {ocr_filename}")
                            ocr_file_id = await self.retrieval_client.upload_file(ocr_path, original_filename=ocr_filename)
                            
                            if ocr_file_id:
                                logger.info(f"[OCR] Successfully uploaded OCR text with ID: {ocr_file_id}")
                                # Update text preview with OCR content
                                text_preview = ocr_text[:1000]
                            
                            # Clean up OCR temp file
                            if os.path.exists(ocr_path):
                                os.unlink(ocr_path)
                    except Exception as ocr_error:
                        logger.warning(f"[OCR] OCR processing failed: {str(ocr_error)}")
                
                # Upload original file
                file_id = await self.retrieval_client.upload_file(file_path, original_filename=original_filename)
                
            except Exception as e:
                logger.warning(f"OpenAI upload failed (will continue): {str(e)}")
            
            # Create document record in Supabase
            document_data = {
                "filename": file_metadata["filename"],
                "storage_path": storage_path,
                "content_type": file_metadata.get("content_type"),
                "size": file_metadata.get("size"),
                "file_id": file_id,  # OpenAI file ID
                "metadata": file_metadata,
                "text_preview": text_preview,
                "user_id": user_id,
                "uploaded_by_email": uploaded_by_email
            }
            
            document_record = await self.supabase.create_document(document_data)
            
            # Clean up temp file if created
            if file_content and filename and 'temp_path' in locals() and temp_path and os.path.exists(temp_path):
                try:
                    os.unlink(temp_path)
                    logger.debug(f"[TEMP FILE] Cleaned up temporary file: {temp_path}")
                except Exception as e:
                    logger.warning(f"[TEMP FILE] Could not delete temp file {temp_path}: {e}")
            
            logger.info(f"Successfully uploaded document: {file_metadata['filename']}")
            
            # Track successful upload
            self.usage_tracker.track_document_upload(
                doc_id=document_record["document_id"],
                file_size=file_metadata.get("size", 0),
                duration_ms=0,  # TODO: Add proper timing
                success=True
            )
            
            return {
                "status": "success",
                "doc_id": document_record["document_id"],
                "file_id": file_id,
                "metadata": file_metadata,
                "storage_url": await self._get_download_url(document_record["document_id"])
            }
            
        except Exception as e:
            logger.error(f"Error uploading document: {str(e)}")
            
            # Track failed upload
            self.usage_tracker.track_document_upload(
                doc_id="unknown",
                file_size=file_metadata.get("size", 0) if 'file_metadata' in locals() else 0,
                duration_ms=0,
                success=False,
                error=str(e)
            )
            
            # Clean up on error
            if file_content and filename and 'temp_path' in locals() and temp_path and os.path.exists(temp_path):
                try:
                    os.unlink(temp_path)
                    logger.debug(f"[TEMP FILE] Cleaned up temporary file after error: {temp_path}")
                except Exception as cleanup_error:
                    logger.warning(f"[TEMP FILE] Could not delete temp file {temp_path}: {cleanup_error}")
            return {
                "status": "error",
                "error": str(e)
            }
    
    async def download_document(self, doc_id: str) -> Optional[Dict[str, Any]]:
        """Download a document from Supabase"""
        try:
            # Get document record
            document = await self.supabase.get_document(doc_id)
            if not document:
                return None
            
            # Get download URL
            download_url = self.supabase.get_file_url(
                self.STORAGE_BUCKET,
                document["storage_path"],
                expires_in=3600  # 1 hour expiry
            )
            
            return {
                "filename": document["filename"],
                "url": download_url,
                "content_type": document.get("content_type"),
                "size": document.get("size_bytes")
            }
            
        except Exception as e:
            logger.error(f"Error downloading document: {str(e)}")
            return None
    
    async def _get_download_url(self, doc_id: str) -> Optional[str]:
        """Get a signed download URL for a document"""
        try:
            document = await self.supabase.get_document(doc_id)
            if not document or not document.get("storage_path"):
                return None
            
            # Try to get signed URL with shorter timeout
            try:
                return self.supabase.get_file_url(
                    self.STORAGE_BUCKET,
                    document["storage_path"],
                    expires_in=3600
                )
            except Exception as url_error:
                # If signed URL fails, return a placeholder or public URL
                logger.warning(f"Could not create signed URL for {doc_id}: {str(url_error)}")
                # Return None so documents still load without download URLs
                return None
        except Exception as e:
            logger.error(f"Error getting download URL: {str(e)}")
            return None
    
    async def delete_document(self, doc_id: str) -> Dict[str, Any]:
        """Delete a document from the system"""
        try:
            # Get document info
            document = await self.supabase.get_document(doc_id)
            if not document:
                return {
                    "status": "error",
                    "error": "Document not found"
                }
            
            # Delete from OpenAI if file_id exists
            if document.get("openai_file_id"):
                try:
                    await self.retrieval_client.delete_file(document["openai_file_id"])
                    logger.info(f"Deleted from OpenAI: {document['openai_file_id']}")
                except Exception as e:
                    logger.warning(f"Failed to delete from OpenAI: {str(e)}")
            
            # DO NOT delete from Supabase Storage - keep for historical purposes
            # The file remains in the bucket even after document is marked as deleted
            logger.info(f"Keeping file in storage for history: {document.get('storage_path')}")
            
            # Soft delete document record (marks as deleted in database)
            await self.supabase.delete_document(doc_id, soft_delete=True)
            
            logger.info(f"Document marked as deleted: {doc_id} (file kept in storage)")
            
            return {
                "status": "success",
                "doc_id": doc_id,
                "message": "Document deleted from active use but preserved in storage"
            }
            
        except Exception as e:
            logger.error(f"Error deleting document: {str(e)}")
            return {
                "status": "error",
                "error": str(e)
            }
    
    async def list_documents(self) -> List[Dict[str, Any]]:
        """List all documents (Admin view with full details)"""
        try:
            logger.info(f"[SUPABASE DB] Listing all documents")
            
            result = self.supabase.client.table("documents").select("*").eq("status", "active").execute()
            
            if result.data:
                logger.info(f"[SUPABASE DB] Found {len(result.data)} documents")
                return result.data
            else:
                logger.info(f"[SUPABASE DB] No documents found")
                return []
                
        except Exception as e:
            logger.error(f"[SUPABASE DB] Error listing documents: {str(e)}")
            return []
    
    async def list_documents_for_users(self) -> List[Dict[str, Any]]:
        """List documents for regular users (limited information)"""
        try:
            logger.info(f"[SUPABASE DB] Listing documents for users")
            
            # 일반 사용자에게는 기본 정보만 제공 (관리 정보 제외)
            result = self.supabase.client.table("documents").select(
                "document_id, filename, content_type, size_bytes, created_at, text_preview"
            ).eq("status", "active").execute()
            
            if result.data:
                logger.info(f"[SUPABASE DB] Found {len(result.data)} documents for users")
                return result.data
            else:
                logger.info(f"[SUPABASE DB] No documents found for users")
                return []
                
        except Exception as e:
            logger.error(f"[SUPABASE DB] Error listing documents for users: {str(e)}")
            return []
    
    async def get_document(self, doc_id: str) -> Optional[Dict[str, Any]]:
        """Get document information (alias for get_document_details)"""
        return await self.get_document_details(doc_id)
    
    async def get_document_details(self, doc_id: str) -> Optional[Dict[str, Any]]:
        """Get detailed information about a document"""
        document = await self.supabase.get_document(doc_id)
        if document:
            document["download_url"] = await self._get_download_url(doc_id)
        return document
    
    async def search_documents(self, query: str, filters: Optional[Dict[str, Any]] = None) -> List[Dict[str, Any]]:
        """Search documents by metadata"""
        # For now, do a simple client-side search
        # In production, you'd want to use Supabase's full-text search
        documents = await self.supabase.list_documents()
        results = []
        query_lower = query.lower()
        
        for doc in documents:
            # Search in filename and text preview
            filename = doc["filename"].lower()
            preview = (doc.get("text_preview") or "").lower()
            
            if query_lower in filename or query_lower in preview:
                # Apply filters if provided
                if filters and not self._match_filters(doc, filters):
                    continue
                
                results.append({
                    "doc_id": doc["document_id"],
                    "filename": doc["filename"],
                    "relevance": self._calculate_relevance(query_lower, filename, preview),
                    "metadata": doc.get("metadata", {}),
                    "download_url": await self._get_download_url(doc["document_id"])
                })
        
        # Sort by relevance
        results.sort(key=lambda x: x["relevance"], reverse=True)
        
        return results
    
    def _validate_file(self, file_path: str) -> Dict[str, Any]:
        """Validate file before upload"""
        path = Path(file_path)
        
        # Check if file exists
        if not path.exists():
            return {"valid": False, "error": "File not found"}
        
        # Check file extension
        if path.suffix.lower() not in self.SUPPORTED_EXTENSIONS:
            return {
                "valid": False, 
                "error": f"Unsupported file type: {path.suffix}"
            }
        
        # Check file size
        size = path.stat().st_size
        if size > self.MAX_FILE_SIZE:
            return {
                "valid": False,
                "error": f"File too large: {format_file_size(size)} (max: {format_file_size(self.MAX_FILE_SIZE)})"
            }
        
        return {"valid": True}
    
    async def _extract_text_preview(self, file_path: str, preview_length: int = 500) -> str:
        """Extract text preview from document"""
        try:
            ext = Path(file_path).suffix.lower()
            
            if ext == '.pdf':
                return await self._extract_pdf_text(file_path, preview_length)
            elif ext in ['.doc', '.docx']:
                return await self._extract_docx_text(file_path, preview_length)
            elif ext in ['.txt', '.md']:
                return await self._extract_text_file(file_path, preview_length)
            elif ext in ['.ppt', '.pptx']:
                return await self._extract_pptx_text(file_path, preview_length)
            elif ext in ['.xls', '.xlsx', '.csv']:
                return await self._extract_excel_text(file_path, preview_length)
            elif ext in ['.png', '.jpg', '.jpeg', '.tiff', '.bmp']:
                return await self._extract_image_text(file_path, preview_length)
            else:
                return "Preview not available for this file type"
                
        except Exception as e:
            logger.warning(f"Could not extract text preview: {str(e)}")
            return "Preview extraction failed"
    
    async def _extract_pdf_text(self, file_path: str, max_length: int) -> str:
        """Extract text from PDF, with OCR fallback for scanned documents"""
        text = []
        
        # First try standard text extraction
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                
                for page_num in range(min(3, len(pdf_reader.pages))):
                    page = pdf_reader.pages[page_num]
                    page_text = page.extract_text()
                    if page_text:
                        text.append(page_text)
            
            full_text = '\n'.join(text)
            
            # If we got very little text, it might be a scanned PDF
            if len(full_text.strip()) < 50:
                logger.info(f"PDF appears to be scanned. Attempting text extraction...")
                
                # Try Vision API first (more reliable for Korean text)
                try:
                    from .vision_processor import get_vision_processor
                    vision = get_vision_processor()
                    vision_text = await vision.extract_text_from_pdf(file_path, max_pages=3)
                    if vision_text and len(vision_text.strip()) > len(full_text.strip()):
                        full_text = vision_text
                        logger.info(f"Successfully extracted text using Vision API")
                except Exception as vision_error:
                    logger.warning(f"Vision API processing failed: {str(vision_error)}")
                    
                    # Fallback to OCR if available
                    try:
                        from .ocr_processor import get_ocr_processor
                        ocr = get_ocr_processor()
                        
                        if ocr.is_ocr_available():
                            ocr_text = await ocr.extract_text_from_pdf(file_path)
                            if ocr_text and len(ocr_text.strip()) > len(full_text.strip()):
                                full_text = ocr_text
                                logger.info(f"Successfully extracted text using OCR")
                        else:
                            logger.info("Tesseract OCR not available. Text extraction limited.")
                    except Exception as ocr_error:
                        logger.warning(f"OCR processing failed: {str(ocr_error)}")
            
            return full_text[:max_length] + "..." if len(full_text) > max_length else full_text
            
        except Exception as e:
            logger.error(f"Error extracting PDF text: {str(e)}")
            return ""
    
    async def _extract_docx_text(self, file_path: str, max_length: int) -> str:
        """Extract text from Word document"""
        doc = Document(file_path)
        text = []
        
        # Try paragraphs first
        for paragraph in doc.paragraphs[:20]:  # First 20 paragraphs
            if paragraph.text.strip():
                text.append(paragraph.text)
        
        # If no paragraphs, try tables
        if not text and doc.tables:
            for table in doc.tables[:5]:  # First 5 tables
                for row in table.rows[:10]:  # First 10 rows
                    row_text = ' | '.join(cell.text.strip() for cell in row.cells)
                    if row_text.strip():
                        text.append(row_text)
        
        full_text = '\n'.join(text)
        return full_text[:max_length] + "..." if len(full_text) > max_length else full_text
    
    async def _extract_text_file(self, file_path: str, max_length: int) -> str:
        """Extract text from plain text file"""
        async with aiofiles.open(file_path, 'r', encoding='utf-8') as file:
            content = await file.read(max_length + 100)
        
        return content[:max_length] + "..." if len(content) > max_length else content
    
    async def _extract_pptx_text(self, file_path: str, max_length: int) -> str:
        """Extract text from PowerPoint"""
        prs = Presentation(file_path)
        text = []
        
        for slide in prs.slides[:5]:  # First 5 slides
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        
        full_text = '\n'.join(text)
        return full_text[:max_length] + "..." if len(full_text) > max_length else full_text
    
    async def _extract_image_text(self, file_path: str, max_length: int) -> str:
        """Extract text from image using Vision API or OCR"""
        try:
            # Try Vision API first (better for complex layouts and Korean text)
            try:
                from .vision_processor import get_vision_processor
                vision = get_vision_processor()
                vision_text = await vision.extract_text_from_image(file_path)
                if vision_text:
                    logger.info(f"Successfully extracted text from image using Vision API")
                    return vision_text[:max_length] + "..." if len(vision_text) > max_length else vision_text
            except Exception as vision_error:
                logger.warning(f"Vision API failed: {str(vision_error)}")
            
            # Fallback to OCR if available
            from .ocr_processor import get_ocr_processor
            ocr = get_ocr_processor()
            
            if ocr.is_ocr_available():
                ocr_text = await ocr.extract_text_from_image(file_path)
                if ocr_text:
                    logger.info(f"Successfully extracted text from image using OCR")
                    return ocr_text[:max_length] + "..." if len(ocr_text) > max_length else ocr_text
                else:
                    return "No text found in image"
            else:
                logger.info("Neither Vision API nor OCR available for image text extraction")
                return "Text extraction not available"
        except Exception as e:
            logger.error(f"Error extracting text from image: {str(e)}")
            return "Image text extraction failed"
    
    async def _extract_excel_text(self, file_path: str, max_length: int) -> str:
        """Extract text from Excel file"""
        wb = openpyxl.load_workbook(file_path, read_only=True)
        text = []
        
        # Get first sheet
        ws = wb.active
        
        # Read first 10 rows
        for row in ws.iter_rows(max_row=10, values_only=True):
            row_text = ' | '.join(str(cell) for cell in row if cell)
            if row_text:
                text.append(row_text)
        
        wb.close()
        
        full_text = '\n'.join(text)
        return full_text[:max_length] + "..." if len(full_text) > max_length else full_text
    
    def _calculate_relevance(self, query: str, filename: str, preview: str) -> float:
        """Calculate simple relevance score"""
        score = 0.0
        
        # Exact match in filename
        if query in filename:
            score += 2.0
        
        # Word match in filename
        query_words = query.split()
        filename_words = filename.split()
        for word in query_words:
            if word in filename_words:
                score += 1.0
        
        # Match in preview
        if query in preview:
            score += 0.5
        
        return score
    
    def _match_filters(self, doc: Dict[str, Any], filters: Dict[str, Any]) -> bool:
        """Check if document matches filters"""
        for key, value in filters.items():
            if key == "file_type":
                doc_type = parse_file_type(doc["filename"])
                if doc_type != value:
                    return False
            elif key == "min_size":
                if doc.get("size_bytes", 0) < value:
                    return False
            elif key == "max_size":
                if doc.get("size_bytes", 0) > value:
                    return False
            elif key == "uploaded_after":
                if doc["created_at"] < value:
                    return False
            elif key == "uploaded_before":
                if doc["created_at"] > value:
                    return False
        
        return True
    
    async def health_check(self) -> Dict[str, Any]:
        """Check health of document manager"""
        try:
            # Check Supabase connection
            supabase_health = await self.supabase.health_check()
            
            # Count documents
            documents = await self.supabase.list_documents()
            active_docs = len([d for d in documents if d.get("status") == "active"])
            
            return {
                "healthy": supabase_health["healthy"],
                "service": "document_manager",
                "supabase_health": supabase_health,
                "total_documents": len(documents),
                "active_documents": active_docs
            }
        except Exception as e:
            return {
                "healthy": False,
                "service": "document_manager",
                "error": str(e)
            }